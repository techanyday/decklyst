import os
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv
import requests
from io import BytesIO
from flask import url_for
import json
import logging
import re

load_dotenv()
openai.api_key = os.getenv('OPENAI_API_KEY')

# Set up logging
logging.basicConfig(level=logging.INFO)

COLOR_THEMES = {
    'blue': RGBColor(0, 112, 192),
    'green': RGBColor(0, 176, 80),
    'red': RGBColor(192, 0, 0),
    'gray': RGBColor(128, 128, 128),
    'blue_white_gradient': (RGBColor(0, 112, 192), RGBColor(255, 255, 255)),
    'red_gray_gradient': (RGBColor(192, 0, 0), RGBColor(128, 128, 128)),
}

def generate_slides_content(topic, num_slides):
    prompt = f"""
    Generate {num_slides} PowerPoint slides about '{topic}'.
    For each slide, respond in JSON as an array of objects with:
    - title
    - bullets (3-5 bullet points)
    - image_prompt (short description for DALL·E 3)
    """
    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1200,
            temperature=0.7
        )
        slides = json.loads(response.choices[0].message.content)
        return slides
    except openai.RateLimitError:
        logging.error("OpenAI API rate limit exceeded or insufficient quota")
        raise Exception("API quota exceeded. Please try again later or check your billing status.")
    except openai.APIError as e:
        logging.error(f"OpenAI API error: {str(e)}")
        raise Exception("Unable to generate slides content. Please try again later.")
    except Exception as e:
        logging.error(f"Unexpected error in generate_slides_content: {str(e)}")
        raise Exception("An unexpected error occurred. Please try again later.")

def generate_slide_images(slides, user_tier):
    images = []
    for slide in slides:
        prompt = slide['image_prompt']
        try:
            dalle_response = openai.images.generate(
                model="dall-e-3",
                prompt=prompt,
                n=1,
                size="1024x1024"
            )
            image_url = dalle_response.data[0].url
            logging.info(f"DALL·E image URL for '{slide['title']}': {image_url}")
            img_data = requests.get(image_url).content
            img_path = os.path.join('static', 'presentations', f"{slide['title'][:10]}_img.png")
            with open(img_path, 'wb') as f:
                f.write(img_data)
            images.append(img_path)
        except openai.RateLimitError:
            logging.error(f"OpenAI API rate limit exceeded for image generation: {slide['title']}")
            images.append(None)
        except Exception as e:
            logging.error(f"Failed to generate or save image for slide '{slide['title']}': {e}")
            images.append(None)
    return images

def create_pptx(slides, images, color_theme, output_path, user_tier):
    prs = Presentation()
    theme = COLOR_THEMES.get(color_theme, RGBColor(0, 112, 192))
    
    # Slide dimensions and spacing constants (in inches)
    SLIDE_WIDTH = 10
    SLIDE_HEIGHT = 7.5
    MARGIN = 0.5
    TITLE_HEIGHT = 0.8
    TITLE_TOP = MARGIN
    MAX_IMAGE_HEIGHT = 3.5
    MIN_IMAGE_HEIGHT = 2.5
    BULLET_SPACING = 0.25
    
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Set background color or gradient
        fill = slide.background.fill
        if color_theme.endswith('_gradient'):
            fill.gradient()
            fill.gradient_stops[0].color.rgb = theme[0]
            fill.gradient_stops[1].color.rgb = theme[1]
        else:
            fill.solid()
            fill.fore_color.rgb = theme

        # Title placement - always at the top
        title_text = slide_data['title']
        title_text = re.sub(r"^Slide\s*\d+[:\-\.]?\s*", "", title_text)
        title_box = slide.shapes.add_textbox(
            Inches(MARGIN), 
            Inches(TITLE_TOP),
            Inches(SLIDE_WIDTH - 2*MARGIN), 
            Inches(TITLE_HEIGHT)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Calculate content area
        content_top = TITLE_TOP + TITLE_HEIGHT
        content_height = SLIDE_HEIGHT - content_top - MARGIN
        
        # Dynamically adjust image height based on number of bullets
        num_bullets = len(slide_data['bullets'])
        if num_bullets <= 3:
            image_height = MAX_IMAGE_HEIGHT
        else:
            # Reduce image height as bullet points increase
            image_height = max(MIN_IMAGE_HEIGHT, MAX_IMAGE_HEIGHT - (num_bullets - 3) * 0.2)
        
        # Image placement - centered below title
        image_top = content_top + 0.2  # Small gap after title
        if images and len(images) > i and images[i]:
            try:
                # Calculate image dimensions maintaining aspect ratio
                img_shape = slide.shapes.add_picture(
                    images[i],
                    Inches(MARGIN),
                    Inches(image_top),
                    Inches(SLIDE_WIDTH - 2*MARGIN),
                    Inches(image_height)
                )
                
                # Maintain aspect ratio
                aspect_ratio = img_shape.height / img_shape.width
                new_width = min(SLIDE_WIDTH - 2*MARGIN, image_height / aspect_ratio)
                new_height = new_width * aspect_ratio
                
                # Center image horizontally
                x_offset = (SLIDE_WIDTH - new_width) / 2
                img_shape.left = Inches(x_offset)
                img_shape.width = Inches(new_width)
                img_shape.height = Inches(new_height)
                
                actual_image_height = new_height
            except Exception as e:
                logging.error(f"Failed to add image to slide '{title_text}': {e}")
                actual_image_height = 0
        else:
            actual_image_height = 0
        
        # Bullet points placement - below image
        bullets_top = image_top + actual_image_height + 0.3
        remaining_height = SLIDE_HEIGHT - bullets_top - MARGIN
        
        bullet_box = slide.shapes.add_textbox(
            Inches(MARGIN),
            Inches(bullets_top),
            Inches(SLIDE_WIDTH - 2*MARGIN),
            Inches(remaining_height)
        )
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Dynamically adjust font size based on content length and number of bullets
        total_text_length = sum(len(bullet) for bullet in slide_data['bullets'])
        avg_text_length = total_text_length / len(slide_data['bullets'])
        
        if avg_text_length > 100:
            base_font_size = 16
        elif avg_text_length > 50:
            base_font_size = 18
        else:
            base_font_size = 20
            
        bullet_font_size = min(base_font_size, 20 if num_bullets <= 4 else max(14, int(20 - (num_bullets - 4))))
        
        for idx, bullet in enumerate(slide_data['bullets']):
            p = bullet_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(bullet_font_size)
            p.level = 0
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            
            # Adjust alignment based on number of bullets and text length
            if len(slide_data['bullets']) <= 3 and avg_text_length < 50:
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
        
        # Watermark for free users (at bottom right corner)
        if user_tier == 'free':
            watermark_box = slide.shapes.add_textbox(
                Inches(SLIDE_WIDTH - 3),
                Inches(SLIDE_HEIGHT - 0.4),
                Inches(2.5),
                Inches(0.3)
            )
            watermark_frame = watermark_box.text_frame
            watermark_frame.text = "Generated by Decklyst (Free)"
            watermark_frame.paragraphs[0].font.size = Pt(11)
            watermark_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
            watermark_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    prs.save(output_path)
